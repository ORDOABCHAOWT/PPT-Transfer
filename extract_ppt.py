#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT智能文案提取工具 v4.0 - 终极版
功能：按列从左到右，每列内从上到下提取文案，绝不遗漏
使用方法：python3 extract_ppt.py
"""

from pptx import Presentation
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml.ns import qn
import os
import sys

class SmartPPTExtractor:
    def __init__(self, ppt_path):
        try:
            print(f"📂 正在打开文件: {ppt_path}")
            self.prs = Presentation(ppt_path)
            print(f"✅ 文件打开成功，共 {len(self.prs.slides)} 页")
            self.doc = Document()
            # 设置默认字体为微软雅黑
            self.doc.styles['Normal'].font.name = '微软雅黑'
            self.doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
        except Exception as e:
            print(f"❌ 无法打开PPT文件: {str(e)}")
            raise
        
    def extract_all_texts_aggressive(self, slide):
        """
        超级激进提取：不遗漏任何文本
        """
        text_boxes = []
        processed_texts = set()  # 用于去重

        def clean_text(text):
            """清理文本中的非法XML字符"""
            import re
            # 移除NULL字节和控制字符（保留换行和制表符）
            text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
            return text

        def add_text_box(text, left, top, font_size, font_name='微软雅黑', width=0, height=0):
            """添加文本框，自动去重"""
            if not text or len(text.strip()) == 0:
                return

            text = clean_text(text.strip())
            if not text:
                return

            # 使用位置和文本内容作为唯一标识
            unique_key = f"{int(left)}_{int(top)}_{text[:100]}"

            if unique_key not in processed_texts:
                processed_texts.add(unique_key)
                text_boxes.append({
                    'text': text,
                    'left': left,
                    'top': top,
                    'font_size': font_size,
                    'font_name': font_name,
                    'width': width,
                    'height': height,
                })

        def get_position(shape, parent_left=0, parent_top=0):
            """获取形状位置"""
            try:
                left = parent_left + (shape.left if hasattr(shape, 'left') and shape.left else 0)
                top = parent_top + (shape.top if hasattr(shape, 'top') and shape.top else 0)
                return left, top
            except:
                return parent_left, parent_top

        def extract_from_shape(shape, parent_left=0, parent_top=0):
            """递归提取所有文本 - 绝不遗漏"""
            try:
                left, top = get_position(shape, parent_left, parent_top)

                # 1. 处理组合形状 - 递归处理所有子形状
                if hasattr(shape, 'shape_type') and shape.shape_type == 6:  # msoGroup
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            extract_from_shape(sub_shape, left, top)
                    return

                # 2. 表格 - 优先处理
                if hasattr(shape, 'has_table'):
                    try:
                        if shape.has_table and hasattr(shape, 'table'):
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        cell_top = top + (row_idx * 300000)
                                        cell_left = left + (col_idx * 300000)
                                        add_text_box(cell_text, cell_left, cell_top, 11.0, 0, 0)
                            return
                    except Exception as e:
                        print(f"      表格提取错误: {str(e)}")

                # 3. text_frame - 主要提取方法
                extracted_text = None
                font_size = 12.0
                font_name = '微软雅黑'  # 默认字体

                if hasattr(shape, 'text_frame'):
                    try:
                        text_frame = shape.text_frame
                        if text_frame and hasattr(text_frame, 'paragraphs'):
                            all_text = []
                            for paragraph in text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    all_text.append(para_text)
                                    # 获取字号和字体名称
                                    for run in paragraph.runs:
                                        if hasattr(run, 'font'):
                                            if run.font.size:
                                                font_size = run.font.size.pt
                                            if run.font.name:
                                                font_name = run.font.name
                                            break

                            if all_text:
                                extracted_text = '\n'.join(all_text)
                    except Exception as e:
                        print(f"      text_frame提取错误: {str(e)}")

                # 4. 直接text属性（备用）
                if not extracted_text and hasattr(shape, 'text'):
                    try:
                        direct_text = shape.text.strip()
                        if direct_text:
                            extracted_text = direct_text
                    except Exception as e:
                        print(f"      text属性提取错误: {str(e)}")

                # 5. 如果提取到文本，添加到列表
                if extracted_text:
                    width = shape.width if hasattr(shape, 'width') else 0
                    height = shape.height if hasattr(shape, 'height') else 0
                    add_text_box(extracted_text, left, top, font_size, font_name, width, height)
                    print(f"      ✓ 提取到文本: {extracted_text[:30]}...")
                else:
                    print(f"      ✗ 未提取到文本")

            except Exception as e:
                print(f"      💥 形状处理异常: {str(e)}")
                import traceback
                traceback.print_exc()

        # 遍历所有形状
        try:
            shapes_list = list(slide.shapes)
            print(f"  🔍 幻灯片共有 {len(shapes_list)} 个形状对象")

            for idx, shape in enumerate(shapes_list, 1):
                # 调试：显示每个形状的信息
                shape_info = f"形状{idx}"
                try:
                    if hasattr(shape, 'shape_type'):
                        shape_info += f" 类型:{shape.shape_type}"
                    if hasattr(shape, 'name'):
                        shape_info += f" 名称:{shape.name}"
                except:
                    pass
                print(f"    处理 {shape_info}")

                extract_from_shape(shape)
        except Exception as e:
            print(f"    ⚠️ 提取形状时出错: {str(e)}")

        # 提取幻灯片备注
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame and notes_text_frame.text:
                    notes_text = notes_text_frame.text.strip()
                    if notes_text:
                        add_text_box(f"【备注】{notes_text}", 0, 999999, 11.0)
        except:
            pass

        return text_boxes
    
    def column_based_sort(self, text_boxes):
        """
        列优先排序：从左到右分列，每列内从上到下
        这是最符合PPT布局的阅读顺序
        """
        if not text_boxes:
            return []

        print(f"\n  📊 开始列优先排序...")
        print(f"  原始文本框数量: {len(text_boxes)}")

        # 第一步：按left值排序，识别列
        sorted_by_left = sorted(text_boxes, key=lambda x: x['left'])

        # 第二步：动态识别列
        # 使用聚类算法识别列边界
        columns = []
        COLUMN_TOLERANCE = 500000  # 列间距容差，约500px

        for box in sorted_by_left:
            placed = False

            # 尝试将文本框放入已有的列
            for col in columns:
                # 计算该列的平均left值
                avg_left = sum(b['left'] for b in col) / len(col)

                # 如果文本框的left值与列的平均值接近，归入该列
                if abs(box['left'] - avg_left) < COLUMN_TOLERANCE:
                    col.append(box)
                    placed = True
                    break

            # 如果没有合适的列，创建新列
            if not placed:
                columns.append([box])

        print(f"  ✓ 识别到 {len(columns)} 列")

        # 第三步：每列内按top值（从上到下）排序
        for i, col in enumerate(columns):
            col.sort(key=lambda x: x['top'])
            min_left = min(b['left'] for b in col)
            max_left = max(b['left'] for b in col)
            print(f"    列{i+1}: {len(col)} 个文本框 (Left范围: {int(min_left)} - {int(max_left)})")

        # 第四步：按列的left值排序列（确保从左到右）
        columns.sort(key=lambda col: min(b['left'] for b in col))

        # 第五步：按列顺序合并所有文本框
        sorted_boxes = []
        for col in columns:
            sorted_boxes.extend(col)

        print(f"  ✓ 排序完成：共 {len(sorted_boxes)} 个文本框")

        return sorted_boxes
    
    def set_font(self, run, font_name='微软雅黑'):
        """设置字体，支持多种字体回退"""
        try:
            # 尝试设置指定字体
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
        except Exception as e:
            # 如果失败，使用默认字体
            try:
                run.font.name = '微软雅黑'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
            except:
                pass  # 忽略字体设置错误
    
    def export_to_word(self, output_path):
        """导出到Word文档"""
        print(f"\n📄 开始处理PPT文件...\n")
        
        total_text_count = 0
        
        for slide_num, slide in enumerate(self.prs.slides, 1):
            print(f"{'='*70}")
            print(f"处理第 {slide_num}/{len(self.prs.slides)} 页")
            print(f"{'='*70}")
            
            try:
                # 添加幻灯片标题
                heading = self.doc.add_heading(f'幻灯片 {slide_num}', level=1)
                heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                for run in heading.runs:
                    self.set_font(run)
                
                # 激进式提取所有文本
                text_boxes = self.extract_all_texts_aggressive(slide)
                print(f"  ✓ 提取到 {len(text_boxes)} 个文本框")
                
                if not text_boxes:
                    print(f"  ⚠️  该页没有文本内容")
                    para = self.doc.add_paragraph("【此页无文本内容】")
                    self.set_font(para.runs[0])
                    self.doc.add_page_break()
                    continue
                
                # 按列优先排序
                sorted_boxes = self.column_based_sort(text_boxes)
                
                print(f"\n  📝 提取文本详细信息（共{len(sorted_boxes)}条）:")

                # 写入Word并显示详细调试信息
                for idx, tb in enumerate(sorted_boxes, 1):
                    text = tb['text']
                    font_size = tb['font_size']
                    font_name = tb.get('font_name', '微软雅黑')  # 获取原始字体名称

                    # 显示提取的文本预览（带详细位置和字体）
                    preview = text.replace('\n', ' ')[:50] + "..." if len(text) > 50 else text.replace('\n', ' ')
                    print(f"  [{idx:2d}] Left:{int(tb['left']):7d} Top:{int(tb['top']):7d} Size:{font_size:4.1f}pt Font:{font_name} | {preview}")

                    # 根据字号判断样式
                    if font_size >= 22:
                        # 大标题
                        para = self.doc.add_heading(text, level=2)
                        for run in para.runs:
                            self.set_font(run, font_name)
                    elif font_size >= 16:
                        # 中标题
                        para = self.doc.add_paragraph()
                        run = para.add_run(text)
                        run.font.size = Pt(15)
                        run.font.bold = True
                        self.set_font(run, font_name)
                    elif font_size >= 13:
                        # 小标题
                        para = self.doc.add_paragraph()
                        run = para.add_run(text)
                        run.font.size = Pt(12)
                        self.set_font(run, font_name)
                    else:
                        # 正文
                        para = self.doc.add_paragraph(text)
                        if para.runs:
                            para.runs[0].font.size = Pt(11)
                            self.set_font(para.runs[0], font_name)
                
                total_text_count += len(sorted_boxes)
                
                # 幻灯片之间添加分隔
                self.doc.add_page_break()
                
            except Exception as e:
                print(f"❌ 处理第 {slide_num} 页时出错: {str(e)}")
                import traceback
                traceback.print_exc()
                continue
        
        # 保存文档
        try:
            self.doc.save(output_path)
            print(f"\n{'='*70}")
            print(f"✅ 导出成功!")
            print(f"{'='*70}")
            print(f"📊 统计信息:")
            print(f"   - 总页数: {len(self.prs.slides)}")
            print(f"   - 提取文本块: {total_text_count}")
            print(f"   - 字体: 微软雅黑")
            print(f"   - 输出文件: {output_path}")
            print(f"{'='*70}")
        except Exception as e:
            print(f"❌ 保存Word文档时出错: {str(e)}")
            raise

    def export_to_word_with_progress(self, output_path, progress_callback=None):
        """导出到Word文档，支持进度回调"""
        total_text_count = 0
        total_slides = len(self.prs.slides)

        for slide_num, slide in enumerate(self.prs.slides, 1):
            if progress_callback:
                progress_callback(slide_num, total_slides, f'处理第 {slide_num}/{total_slides} 页...')

            try:
                # 添加幻灯片标题
                heading = self.doc.add_heading(f'幻灯片 {slide_num}', level=1)
                heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                for run in heading.runs:
                    self.set_font(run)

                # 激进式提取所有文本
                text_boxes = self.extract_all_texts_aggressive(slide)

                if not text_boxes:
                    para = self.doc.add_paragraph("【此页无文本内容】")
                    self.set_font(para.runs[0])
                    self.doc.add_page_break()
                    continue

                # 按列优先排序
                sorted_boxes = self.column_based_sort(text_boxes)

                # 写入Word
                for tb in sorted_boxes:
                    text = tb['text']
                    font_size = tb['font_size']
                    font_name = tb.get('font_name', '微软雅黑')  # 获取原始字体名称

                    # 根据字号判断样式
                    if font_size >= 22:
                        para = self.doc.add_heading(text, level=2)
                        for run in para.runs:
                            self.set_font(run, font_name)
                    elif font_size >= 16:
                        para = self.doc.add_paragraph()
                        run = para.add_run(text)
                        run.font.size = Pt(15)
                        run.font.bold = True
                        self.set_font(run, font_name)
                    elif font_size >= 13:
                        para = self.doc.add_paragraph()
                        run = para.add_run(text)
                        run.font.size = Pt(12)
                        self.set_font(run, font_name)
                    else:
                        para = self.doc.add_paragraph(text)
                        if para.runs:
                            para.runs[0].font.size = Pt(11)
                            self.set_font(para.runs[0], font_name)

                total_text_count += len(sorted_boxes)

                # 幻灯片之间添加分隔
                self.doc.add_page_break()

            except Exception as e:
                if progress_callback:
                    progress_callback(slide_num, total_slides, f'处理第 {slide_num} 页时出错: {str(e)}')
                continue

        # 保存文档
        self.doc.save(output_path)
        if progress_callback:
            progress_callback(total_slides, total_slides, '导出完成！')
        return total_text_count


def select_ppt_file():
    """让用户选择PPT文件"""
    try:
        import tkinter as tk
        from tkinter import filedialog
        
        root = tk.Tk()
        root.withdraw()
        
        file_path = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=[
                ("PowerPoint文件", "*.pptx"),
                ("所有文件", "*.*")
            ]
        )
        
        return file_path
    except:
        return None


def main():
    print("\n" + "="*70)
    print("PPT智能文案提取工具 v4.0 - 终极完全版")
    print("特性: 列优先排序 | 零遗漏提取 | 微软雅黑字体")
    print("="*70 + "\n")
    
    # 获取PPT文件路径
    ppt_path = None
    
    if len(sys.argv) > 1:
        ppt_path = sys.argv[1]
    
    if not ppt_path:
        try:
            print("📂 正在打开文件选择对话框...")
            ppt_path = select_ppt_file()
        except:
            pass
    
    if not ppt_path:
        print("\n请输入PPT文件路径（可以直接拖拽文件到终端）:")
        ppt_path = input("路径: ").strip().strip("'\"")
    
    if not ppt_path:
        print("❌ 未指定文件，程序退出")
        return
    
    if not os.path.exists(ppt_path):
        print(f"❌ 文件不存在: {ppt_path}")
        return
    
    # 生成输出路径
    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    output_dir = os.path.dirname(ppt_path) or os.path.expanduser("~/Desktop")
    output_path = os.path.join(output_dir, f"{base_name}_完整提取.docx")
    
    # 执行提取
    try:
        extractor = SmartPPTExtractor(ppt_path)
        extractor.export_to_word(output_path)

        print("\n✨ 完成！按Enter键退出...")
        try:
            input()
        except EOFError:
            pass  # 非交互模式下忽略

    except KeyboardInterrupt:
        print("\n\n⚠️  用户中断操作")
    except Exception as e:
        print(f"\n❌ 发生错误: {str(e)}")
        import traceback
        traceback.print_exc()
        print("\n按Enter键退出...")
        try:
            input()
        except EOFError:
            pass  # 非交互模式下忽略


if __name__ == "__main__":
    main()
